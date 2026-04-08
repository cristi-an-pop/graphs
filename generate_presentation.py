from pathlib import Path

import matplotlib.pyplot as plt
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt


def graph_metrics(graph: nx.Graph) -> dict:
    degrees = [d for _, d in graph.degree()]
    ccs = list(nx.connected_components(graph))
    largest_cc_nodes = max(ccs, key=len)
    gcc = graph.subgraph(largest_cc_nodes).copy()

    return {
        "n": graph.number_of_nodes(),
        "m": graph.number_of_edges(),
        "min_deg": min(degrees),
        "max_deg": max(degrees),
        "avg_deg": sum(degrees) / len(degrees),
        "components": len(ccs),
        "lcc_size": len(largest_cc_nodes),
        "diameter": nx.diameter(gcc),
        "avg_shortest_path": nx.average_shortest_path_length(gcc),
        "avg_clustering": nx.average_clustering(graph),
    }


def make_topology_figure(graphs: list[tuple[str, nx.Graph]], out_path: Path) -> None:
    fig, axes = plt.subplots(1, 3, figsize=(18, 6))

    for ax, (name, g) in zip(axes, graphs):
        pos = nx.spring_layout(g, seed=42)
        nx.draw(
            g,
            pos,
            ax=ax,
            node_size=14,
            alpha=0.6,
            edge_color="gray",
            width=0.5,
        )
        ax.set_title(f"{name}\nNodes: {g.number_of_nodes()}", fontsize=10)
        ax.set_axis_off()

    plt.tight_layout()
    fig.savefig(out_path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def make_degree_figure(graphs: list[tuple[str, nx.Graph]], out_path: Path) -> None:
    fig, axes = plt.subplots(1, 3, figsize=(18, 4.5))

    for ax, (name, g) in zip(axes, graphs):
        degree_counts = {}
        for _, degree in g.degree():
            degree_counts[degree] = degree_counts.get(degree, 0) + 1

        x = sorted(degree_counts.keys())
        y = [degree_counts[d] for d in x]

        ax.scatter(x, y, s=12)
        ax.set_xscale("log")
        ax.set_yscale("log")
        ax.set_title(name, fontsize=10)
        ax.set_xlabel("Degree")
        ax.set_ylabel("Count")

    plt.tight_layout()
    fig.savefig(out_path, dpi=180, bbox_inches="tight")
    plt.close(fig)


def add_bullets(slide, title: str, bullets: list[str]) -> None:
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.2))
    tf = textbox.text_frame
    tf.clear()

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)


def main() -> None:
    root = Path(__file__).resolve().parent
    figures_dir = root / "presentation_assets"
    figures_dir.mkdir(exist_ok=True)

    # Load real graph and generate null models with matching size/density.
    g_real = nx.read_edgelist(root / "out.dimacs10-netscience", comments="%", create_using=nx.Graph())
    n = g_real.number_of_nodes()
    m = g_real.number_of_edges()

    max_edges = n * (n - 1) / 2
    p = m / max_edges
    m_ba = max(1, int(round(m / n)))

    g_er = nx.erdos_renyi_graph(n, p, seed=42)
    g_ba = nx.barabasi_albert_graph(n, m_ba, seed=42)

    graphs = [
        ("NetScience (Real)", g_real),
        ("Erdos-Renyi", g_er),
        ("Barabasi-Albert", g_ba),
    ]

    topology_png = figures_dir / "topology_comparison.png"
    degree_png = figures_dir / "degree_distribution.png"

    make_topology_figure(graphs, topology_png)
    make_degree_figure(graphs, degree_png)

    metrics = {name: graph_metrics(g) for name, g in graphs}

    prs = Presentation()

    # Slide 1: title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Network Science Co-authorship Analysis"
    slide.placeholders[1].text = "dimacs10-netscience vs Erdos-Renyi and Barabasi-Albert"

    # Slide 2: dataset
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(
        slide,
        "Dataset Overview",
        [
            "Dataset: dimacs10-netscience (co-authorship network)",
            f"Nodes (authors): {metrics['NetScience (Real)']['n']}",
            f"Edges (co-authorship links): {metrics['NetScience (Real)']['m']}",
            "Graph type: undirected, unweighted, no self-loops",
        ],
    )

    # Slide 3: model setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(
        slide,
        "Experimental Setup",
        [
            "Null model 1: Erdos-Renyi G(n, p) using real-network density",
            "Null model 2: Barabasi-Albert using m ~ M/N",
            "Purpose: compare random, preferential-attachment, and real collaboration structure",
            f"ER p used: {p:.6f}; BA m used: {m_ba}",
        ],
    )

    # Slide 4: topology figure
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Topology Comparison"
    slide.shapes.add_picture(str(topology_png), Inches(0.5), Inches(1.1), width=Inches(12.3))

    # Slide 5: degree distribution figure
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Degree Distribution (log-log)"
    slide.shapes.add_picture(str(degree_png), Inches(0.5), Inches(1.1), width=Inches(12.3))

    # Slide 6: key findings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(
        slide,
        "Key Findings",
        [
            f"Real graph clustering is high: {metrics['NetScience (Real)']['avg_clustering']:.3f}",
            f"ER clustering is very low: {metrics['Erdos-Renyi']['avg_clustering']:.3f}",
            f"BA max degree is much larger than ER: {metrics['Barabasi-Albert']['max_deg']} vs {metrics['Erdos-Renyi']['max_deg']}",
            "Interpretation: real network shows both small-world clustering and hub structure",
        ],
    )

    # Slide 7: summary metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(
        slide,
        "Core Metrics Summary",
        [
            f"Real: components={metrics['NetScience (Real)']['components']}, LCC={metrics['NetScience (Real)']['lcc_size']}, diameter={metrics['NetScience (Real)']['diameter']}",
            f"ER: components={metrics['Erdos-Renyi']['components']}, LCC={metrics['Erdos-Renyi']['lcc_size']}, diameter={metrics['Erdos-Renyi']['diameter']}",
            f"BA: components={metrics['Barabasi-Albert']['components']}, LCC={metrics['Barabasi-Albert']['lcc_size']}, diameter={metrics['Barabasi-Albert']['diameter']}",
            "Conclusion: real collaboration is neither purely random nor purely preferential attachment",
        ],
    )

    output_path = root / "NetworkScience_Presentation.pptx"
    prs.save(output_path)

    print(f"Presentation generated: {output_path}")
    print(f"Assets folder: {figures_dir}")


if __name__ == "__main__":
    main()
